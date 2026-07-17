from pathlib import Path
import json

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "presentacion" / "final"
ASSETS = OUT / "assets_redesign"
ASSETS.mkdir(parents=True, exist_ok=True)
PPTX = OUT / "Presentacion_Final_Churn_Redesign.pptx"

DATA = ROOT / "datos" / "raw" / "WA_Fn-UseC_-Telco-Customer-Churn.csv"
OUTPUTS = ROOT / "modelos" / "resultados"

INK = RGBColor(22, 28, 39)
MUTED = RGBColor(94, 104, 119)
PAPER = RGBColor(249, 250, 252)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(220, 226, 235)
BLUE = RGBColor(35, 86, 194)
TEAL = RGBColor(34, 151, 153)
ORANGE = RGBColor(231, 126, 55)
RED = RGBColor(200, 71, 82)
GREEN = RGBColor(60, 154, 96)
VIOLET = RGBColor(104, 86, 204)

M = {
    "ink": "#161c27",
    "muted": "#5e6877",
    "paper": "#f9fafc",
    "line": "#dce2eb",
    "blue": "#2356c2",
    "teal": "#229799",
    "orange": "#e77e37",
    "red": "#c84752",
    "green": "#3c9a60",
    "violet": "#6856cc",
}


def pct(x, d=1):
    return f"{100*x:.{d}f}%"


def num(x, d=3):
    return f"{x:.{d}f}"


def save(path):
    plt.tight_layout()
    plt.savefig(path, dpi=190, bbox_inches="tight", facecolor="white")
    plt.close()


def load():
    df = pd.read_csv(DATA)
    df["TotalCharges"] = pd.to_numeric(df["TotalCharges"], errors="coerce").fillna(0)
    df["churn_num"] = (df["Churn"] == "Yes").astype(int)
    metrics = json.loads((OUTPUTS / "model_metrics.json").read_text(encoding="utf-8"))
    cv = pd.read_csv(OUTPUTS / "model_selection_cv_results.csv")
    thresholds = pd.read_csv(OUTPUTS / "threshold_search_results.csv").sort_values("threshold")
    preds = pd.read_csv(OUTPUTS / "test_predictions.csv")
    perm = pd.read_csv(OUTPUTS / "feature_importance_permutation_f2.csv")
    shap = pd.read_csv(OUTPUTS / "feature_importance_shap.csv")
    checks = json.loads((OUTPUTS / "leakage_overfit_checks.json").read_text(encoding="utf-8"))
    return df, metrics, cv, thresholds, preds, perm, shap, checks


def style_ax(ax):
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(alpha=0.16)
    ax.tick_params(colors=M["muted"], labelsize=9)
    ax.title.set_color(M["ink"])
    ax.xaxis.label.set_color(M["muted"])
    ax.yaxis.label.set_color(M["muted"])


def make_assets(df, metrics, cv, thresholds, preds, perm, shap):
    # EDA: churn rate vs multiple variables/categories. More legible grouped bars.
    variables = {
        "Contract": "Contrato",
        "InternetService": "Internet",
        "PaymentMethod": "Metodo de pago",
        "OnlineSecurity": "Seguridad online",
        "TechSupport": "Soporte tecnico",
        "PaperlessBilling": "Factura digital",
    }
    rows = []
    for col, label in variables.items():
        rates = df.groupby(col)["churn_num"].mean().sort_values(ascending=False)
        for cat, rate in rates.items():
            rows.append({"variable": label, "category": str(cat), "label": f"{label}: {str(cat)}", "rate": rate})
    plot = pd.DataFrame(rows).sort_values("rate", ascending=True)
    fig, ax = plt.subplots(figsize=(9.5, 5.2))
    colors = [M["orange"] if v >= df["churn_num"].mean() else M["blue"] for v in plot["rate"]]
    ax.barh(plot["label"], plot["rate"], color=colors, height=0.58)
    ax.axvline(df["churn_num"].mean(), color=M["red"], linestyle="--", linewidth=1.5)
    ax.text(df["churn_num"].mean() + 0.008, -0.75, f"media {pct(df['churn_num'].mean())}", color=M["red"], fontsize=9, weight="bold")
    for i, r in enumerate(plot.itertuples()):
        ax.text(r.rate + 0.006, i, pct(r.rate), va="center", fontsize=8.8, color=M["ink"], weight="bold")
    ax.set_xlim(0, 0.52)
    ax.set_xlabel("Tasa de churn")
    ax.set_title("Tasa de churn por categoria: principales variables", fontsize=14, weight="bold")
    style_ax(ax)
    ax.grid(axis="x", alpha=0.18)
    save(ASSETS / "eda_churn_rate_variables.png")

    # Model PR-AUC with 10 tests.
    cvp = cv.sort_values("pr_auc_mean", ascending=True).copy()
    labels = {
        "dummy_most_frequent": "Dummy",
        "gaussian_nb": "Gaussian NB",
        "decision_tree": "Decision Tree",
        "ada_boost": "AdaBoost",
        "extra_trees": "Extra Trees",
        "lightgbm": "LightGBM",
        "random_forest": "Random Forest",
        "logistic_regression": "Logistic Reg.",
        "gradient_boosting": "Gradient Boosting",
        "catboost": "CatBoost",
    }
    cvp["label"] = cvp["model"].map(labels).fillna(cvp["model"])
    fig, ax = plt.subplots(figsize=(8.6, 5.0))
    colors = [M["orange"] if m == metrics["selected_model"] else M["blue"] for m in cvp["model"]]
    ax.barh(cvp["label"], cvp["pr_auc_mean"], xerr=cvp["pr_auc_std"], color=colors, height=0.62)
    ax.axvline(df["churn_num"].mean(), color=M["red"], linestyle="--", linewidth=1.2)
    ax.text(df["churn_num"].mean() + 0.006, -0.45, "no skill", color=M["red"], fontsize=8, weight="bold")
    for i, r in enumerate(cvp.itertuples()):
        ax.text(r.pr_auc_mean + 0.004, i, f"{r.pr_auc_mean:.3f}", va="center", fontsize=8.5, color=M["ink"])
    ax.set_xlim(0.23, max(cvp["pr_auc_mean"]) + 0.045)
    ax.set_xlabel("PR-AUC promedio en validacion cruzada")
    ax.set_title("10 pruebas de modelos: comparacion por PR-AUC", fontsize=14, weight="bold")
    style_ax(ax)
    save(ASSETS / "modelos_10_prauc.png")

    # Compact EDA target + tenure.
    fig, axes = plt.subplots(1, 2, figsize=(9, 3.7))
    churn = df["Churn"].value_counts(normalize=True).reindex(["No", "Yes"])
    axes[0].bar(["No churn", "Churn"], churn.values, color=[M["blue"], M["orange"]], width=0.55)
    axes[0].set_ylim(0, 1)
    axes[0].set_title("Distribucion del target", fontsize=12, weight="bold")
    for i, v in enumerate(churn.values):
        axes[0].text(i, v + 0.025, pct(v), ha="center", fontsize=11, weight="bold", color=M["ink"])
    tg = df.assign(tenure_group=pd.cut(df["tenure"], [-1, 6, 12, 24, 48, np.inf], labels=["0-6", "7-12", "13-24", "25-48", "49+"]))
    rates = tg.groupby("tenure_group", observed=False)["churn_num"].mean()
    axes[1].plot(rates.index.astype(str), rates.values, marker="o", color=M["teal"], linewidth=2.7)
    axes[1].set_title("Churn por antiguedad", fontsize=12, weight="bold")
    axes[1].set_ylabel("Tasa churn")
    for ax in axes:
        style_ax(ax)
    save(ASSETS / "eda_target_tenure.png")

    # Threshold.
    fig, ax = plt.subplots(figsize=(8.2, 4.3))
    ax.plot(thresholds["threshold"], thresholds["recall"], color=M["blue"], linewidth=2.5, label="Recall")
    ax.plot(thresholds["threshold"], thresholds["precision"], color=M["orange"], linewidth=2.5, label="Precision")
    ax.plot(thresholds["threshold"], thresholds["f2"], color=M["green"], linewidth=2.5, label="F2")
    ax.axvline(metrics["decision_threshold"], color=M["red"], linestyle="--")
    ax.set_xlim(0.1, 0.9)
    ax.set_ylim(0, 1.02)
    ax.set_xlabel("Umbral")
    ax.set_title("Umbral operativo: balance precision / recall", fontsize=14, weight="bold")
    ax.legend(frameon=False, loc="upper right")
    style_ax(ax)
    save(ASSETS / "threshold_precision_recall.png")

    # Confusion matrix.
    cm = pd.crosstab(preds["y_true_label"], preds["churn_prediction_label"]).reindex(index=["No", "Yes"], columns=["No", "Yes"], fill_value=0)
    fig, ax = plt.subplots(figsize=(5.0, 4.2))
    im = ax.imshow(cm.values, cmap="Blues")
    ax.set_xticks([0, 1], ["Pred No", "Pred Churn"])
    ax.set_yticks([0, 1], ["Real No", "Real Churn"])
    ax.set_title("Matriz de confusion en test", fontsize=14, weight="bold")
    for i in range(2):
        for j in range(2):
            ax.text(j, i, str(cm.values[i, j]), ha="center", va="center", fontsize=17, weight="bold", color=M["ink"])
    ax.tick_params(labelsize=10)
    plt.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    save(ASSETS / "confusion_matrix.png")

    # Feature importance.
    top = perm.head(10).sort_values("importance_mean")
    fig, ax = plt.subplots(figsize=(7.7, 4.7))
    ax.barh(top["feature"], top["importance_mean"], xerr=top["importance_std"], color=M["teal"], height=0.6)
    ax.set_xlabel("Caida de F2 al permutar")
    ax.set_title("Feature importance por permutation importance", fontsize=14, weight="bold")
    style_ax(ax)
    save(ASSETS / "importance_permutation.png")

    top = shap.head(10).sort_values("mean_abs_shap")
    fig, ax = plt.subplots(figsize=(7.7, 4.7))
    ax.barh(top["feature"], top["mean_abs_shap"], color=M["violet"], height=0.6)
    ax.set_xlabel("mean(|SHAP|)")
    ax.set_title("SHAP global del modelo final", fontsize=14, weight="bold")
    style_ax(ax)
    save(ASSETS / "importance_shap.png")


def add_bg(slide):
    bg = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(5.625))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    band = slide.shapes.add_shape(1, 0, 0, Inches(0.18), Inches(5.625))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()


def textbox(slide, text, x, y, w, h, size=12, color=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def title(slide, section, main):
    textbox(slide, section.upper(), 0.48, 0.24, 3.1, 0.28, 8.5, BLUE, True)
    textbox(slide, main, 0.48, 0.52, 8.9, 0.65, 23, INK, True)


def bullets(slide, items, x, y, w, h, size=12):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(5)
    return box


def card(slide, label, value, note, x, y, w, h, color=BLUE):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    stripe = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.08))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()
    textbox(slide, label, x + 0.15, y + 0.16, w - 0.3, 0.22, 8.2, MUTED, True)
    if value:
        textbox(slide, value, x + 0.15, y + 0.42, w - 0.3, 0.42, 20, INK, True)
    if note:
        textbox(slide, note, x + 0.15, y + (0.86 if value else 0.42), w - 0.3, h - 0.9, 9.2, MUTED)


def image(slide, path, x, y, w, h=None):
    if h:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))


def footer(slide, n):
    textbox(slide, f"Churn Telco · Examen Final · {n:02d}", 0.48, 5.28, 4.5, 0.18, 7.5, MUTED)


def small_note(slide, text, x, y, w, h, color=MUTED):
    return textbox(slide, text, x, y, w, h, 9.5, color)


def build_deck(df, metrics, cv, preds, checks):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    def new(section, main, n):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(s)
        title(s, section, main)
        footer(s, n)
        return s

    raw_predictors = 19
    total_vars = 21
    churn_count = int(df["churn_num"].sum())
    no_churn_count = int(len(df) - churn_count)

    # 1 portada
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, Inches(10), Inches(5.625))
    bg.fill.solid(); bg.fill.fore_color.rgb = INK; bg.line.fill.background()
    textbox(s, "ANALISIS PREDICTIVO - EXAMEN FINAL", 0.6, 0.55, 5.5, 0.3, 9.5, RGBColor(146, 196, 255), True)
    textbox(s, "Prediccion de churn en clientes Telco", 0.6, 1.15, 6.8, 0.95, 31, WHITE, True)
    textbox(s, "Caso de negocio, modelo predictivo y accion comercial sobre clientes con riesgo de baja", 0.64, 2.1, 7.3, 0.35, 13, RGBColor(205, 215, 230))
    card(s, "CLIENTES", f"{len(df):,}".replace(",", "."), "dataset IBM", 0.75, 3.55, 1.7, 0.95, TEAL)
    card(s, "CHURN", pct(df["churn_num"].mean()), f"{churn_count} clientes", 2.75, 3.55, 1.7, 0.95, ORANGE)
    card(s, "MODELO", "CatBoost", "mejor PR-AUC", 4.75, 3.55, 1.7, 0.95, GREEN)
    card(s, "RECALL", pct(metrics["test_recall"]), "test", 6.75, 3.55, 1.7, 0.95, BLUE)

    # 2 caso de negocio
    s = new("1 - Caso de negocio", "Predecir churn para priorizar retencion rentable", 2)
    card(s, "TARGET", "Churn", "Yes si el cliente se dio de baja", 0.65, 1.32, 2.1, 0.95, ORANGE)
    card(s, "BASE", f"{len(df):,}".replace(",", "."), f"{churn_count} churners vs {no_churn_count} no churn", 3.0, 1.32, 2.1, 0.95, TEAL)
    card(s, "TASA", pct(df["churn_num"].mean()), "clase positiva minoritaria", 5.35, 1.32, 1.65, 0.95, RED)
    card(s, "OBJETIVO", "score", "probabilidad de churn por cliente", 7.25, 1.32, 2.0, 0.95, BLUE)
    bullets(s, [
        "Problema: una telco pierde ingreso recurrente cuando un cliente abandona el servicio.",
        "Decision de negocio: ordenar clientes activos por riesgo y priorizar acciones de retencion.",
        "Regla economica: adquirir un cliente nuevo suele costar 5 a 25 veces mas que retener uno existente (HBR/Bain).",
        "Impacto esperado: incluso mejoras chicas de retencion pueden tener efecto grande en rentabilidad; el modelo ayuda a enfocar el esfuerzo comercial."
    ], 0.78, 2.75, 8.5, 1.55, 12.2)
    small_note(s, "Nota: el modelo no decide la oferta; entrega un ranking para que negocio defina accion, costo y capacidad.", 0.85, 4.75, 8.2, 0.28)

    # 3 dataset
    s = new("2 - Dataset", "Variables predictoras agrupadas y limpieza realizada", 3)
    card(s, "TOTAL", str(total_vars), "columnas originales", 0.55, 1.25, 1.55, 0.82, BLUE)
    card(s, "PREDICTORAS", str(raw_predictors), "sin customerID ni Churn", 2.28, 1.25, 1.75, 0.82, GREEN)
    card(s, "TARGET", "1", "Churn", 4.23, 1.25, 1.25, 0.82, ORANGE)
    card(s, "ID", "1", "excluido del modelo", 5.68, 1.25, 1.45, 0.82, RED)
    card(s, "LIMPIEZA", "11", "TotalCharges vacio", 7.35, 1.25, 1.9, 0.82, TEAL)
    rows = [
        "Demograficas (4): genero, senior, pareja, dependientes. Senal baja/moderada; dependientes y pareja reducen churn observado.",
        "Servicios (9): telefono, internet, seguridad, backup, proteccion, soporte, streaming. Fibra y falta de soporte/seguridad elevan riesgo.",
        "Contrato (2): tenure y tipo de contrato. Contrato mensual y baja antiguedad son los drivers mas fuertes.",
        "Facturacion (4): paperless, metodo de pago, MonthlyCharges, TotalCharges. Pago electronico y ciertos niveles de cargo aportan senal."
    ]
    bullets(s, rows, 0.75, 2.45, 8.6, 1.9, 11.3)
    small_note(s, "Limpieza: TotalCharges venia como texto; los 11 vacios correspondian a clientes nuevos con tenure = 0 y se imputaron como 0. Se agregaron features derivadas reproducibles.", 0.78, 4.72, 8.5, 0.35)

    # 4 EDA
    s = new("3 - EDA", "La tasa de churn varia fuerte por categoria", 4)
    image(s, ASSETS / "eda_churn_rate_variables.png", 0.55, 1.18, 8.9)
    small_note(s, "Lectura: contrato mensual, fibra optica y ausencia de soporte/seguridad se ubican por encima de la media de churn. Esto anticipa variables importantes del modelo.", 0.75, 4.88, 8.45, 0.3)

    # 5 evaluacion metricas
    s = new("4 - Evaluacion", "Por que PR-AUC para modelo y F2 para umbral", 5)
    card(s, "DESBALANCE", pct(df["churn_num"].mean()), "churn positivo", 0.65, 1.35, 1.9, 0.95, RED)
    card(s, "NO SKILL", num(df["churn_num"].mean()), "PR-AUC base", 2.85, 1.35, 1.8, 0.95, ORANGE)
    card(s, "PR-AUC", "ranking", "compara modelos", 4.95, 1.35, 1.8, 0.95, BLUE)
    card(s, "F2", "umbral", "prioriza recall", 7.05, 1.35, 1.8, 0.95, GREEN)
    bullets(s, [
        "Accuracy engana: predecir siempre No Churn acierta 73,5% pero no detecta ningun churner.",
        "PR-AUC evalua si el modelo pone arriba del ranking a los clientes que realmente churnean.",
        "F2 se usa despues para elegir el punto de corte: pesa mas recall, pero precision sigue entrando en la formula.",
        "Separar ranking y umbral evita mezclar dos decisiones distintas."
    ], 0.9, 2.8, 8.1, 1.55, 12.5)

    # 6 validacion
    s = new("4 - Evaluacion", "Validacion: todo se decide dentro de train", 6)
    card(s, "TRAIN", "5.634", "80%, estratificado", 0.75, 1.35, 1.7, 0.95, BLUE)
    card(s, "TEST", "1.409", "20%, cerrado", 2.75, 1.35, 1.7, 0.95, TEAL)
    card(s, "CV", "3 folds", "solo train", 4.75, 1.35, 1.7, 0.95, GREEN)
    card(s, "SEED", "42", "reproducible", 6.75, 1.35, 1.7, 0.95, VIOLET)
    bullets(s, [
        "Split estratificado: la proporcion de churn queda estable en train y test.",
        "GridSearchCV: cada modelo se compara con la misma validacion cruzada y la misma metrica PR-AUC.",
        "Pipeline: imputacion, escalado y one-hot se ajustan dentro de cada fold, evitando leakage de preprocesamiento.",
        "El umbral se elige en validacion interna del train; el test se usa una sola vez para reportar performance final."
    ], 0.95, 2.85, 8.0, 1.45, 12.3)

    # 7 baseline
    s = new("5 - Baseline", "Baseline paso a paso: Dummy y regresion logistica", 7)
    cvi = cv.set_index("model")
    card(s, "Dummy", num(cvi.loc["dummy_most_frequent", "pr_auc_mean"]), "predice clase mayoritaria", 0.7, 1.35, 2.1, 0.95, RED)
    card(s, "Logistic", num(cvi.loc["logistic_regression", "pr_auc_mean"]), "PR-AUC CV", 3.05, 1.35, 2.1, 0.95, BLUE)
    card(s, "Predictoras", str(raw_predictors), "variables originales + FE", 5.4, 1.35, 2.1, 0.95, GREEN)
    card(s, "Logistic F2", num(cvi.loc["logistic_regression", "f2_mean"]), "baseline fuerte", 7.75, 1.35, 1.55, 0.95, ORANGE)
    bullets(s, [
        "1. Dummy sirve como piso: si un modelo no supera la prevalencia en PR-AUC, no aporta informacion.",
        "2. Regresion logistica es baseline real: simple, interpretable y rapida de entrenar.",
        "3. Usa el mismo pipeline y las mismas variables que los modelos complejos, por eso la comparacion es justa.",
        "4. Queda cerca del ganador, lo que muestra que el problema tiene senal clara y no depende solo de modelos sofisticados."
    ], 0.85, 2.75, 8.4, 1.7, 12.2)

    # 8 seleccion
    s = new("6 - Seleccion de modelos", "10 pruebas comparadas por PR-AUC", 8)
    image(s, ASSETS / "modelos_10_prauc.png", 0.55, 1.12, 5.85)
    bullets(s, [
        "CatBoost queda primero por PR-AUC, pero la diferencia con Gradient Boosting y Logistic es chica.",
        "Gradient Boosting clasico logra mayor precision, pero menor recall: es mas conservador.",
        "Logistic Regression confirma una senal lineal fuerte y funciona como alternativa simple de defensa.",
        "LightGBM fue probado; en esta grilla no supera a CatBoost ni a los baselines fuertes."
    ], 6.55, 1.38, 3.0, 2.65, 11.2)
    small_note(s, "Conclusion: el ganador no se elige por marketing del algoritmo, sino por PR-AUC bajo la misma validacion.", 6.55, 4.45, 3.0, 0.45)

    # 9 modelo final
    s = new("7 - Modelo final", "CatBoost: buen ranking y umbral orientado a detectar churn", 9)
    card(s, "PR-AUC", num(metrics["test_pr_auc"]), "test", 0.6, 1.25, 1.55, 0.85, ORANGE)
    card(s, "Recall", pct(metrics["test_recall"]), "test", 2.35, 1.25, 1.55, 0.85, GREEN)
    card(s, "Precision", pct(metrics["test_precision"]), "test", 4.1, 1.25, 1.55, 0.85, RED)
    card(s, "F2", num(metrics["test_f2"]), "umbral 0.31", 5.85, 1.25, 1.55, 0.85, BLUE)
    card(s, "Accionados", pct(metrics["positive_rate_predicted_test"]), "del test", 7.6, 1.25, 1.55, 0.85, VIOLET)
    image(s, ASSETS / "threshold_precision_recall.png", 0.62, 2.45, 4.85)
    image(s, ASSETS / "confusion_matrix.png", 5.95, 2.25, 3.05)
    small_note(s, "El umbral 0.31 prioriza recall: captura gran parte de los churners, a costa de falsos positivos. En negocio real, el corte deberia ajustarse por costo/capacidad.", 0.82, 4.88, 8.3, 0.32)

    # 10 interpretacion
    s = new("7 - Modelo final", "Interpretacion: consistente con el EDA", 10)
    image(s, ASSETS / "importance_permutation.png", 0.55, 1.18, 4.65)
    image(s, ASSETS / "importance_shap.png", 5.05, 1.18, 4.55)
    textbox(s, "Conclusion: contrato mensual, baja antiguedad, cargos acumulados y senales de soporte/seguridad vuelven a aparecer como drivers predictivos. Esto coincide con el EDA, pero no implica causalidad.", 0.85, 4.78, 8.25, 0.48, 10.5, MUTED, align=PP_ALIGN.CENTER)

    # 11 limitaciones y mejoras
    s = new("8 - Limitaciones / mejoras", "Que falta para llevarlo a una decision productiva", 11)
    textbox(s, "Limitaciones", 0.75, 1.35, 3.8, 0.3, 16, INK, True)
    bullets(s, [
        "Dataset sin dimension temporal: no simula pasado a futuro.",
        "No hay costos reales de contacto, oferta ni perdida de cliente.",
        "Precision moderada: muchas acciones comerciales no terminarian en churn real.",
        "Variables redundantes entre originales y derivadas pueden repartir importancia.",
        "No hay datos de reclamos, uso, satisfaccion ni competencia."
    ], 0.75, 1.82, 4.0, 2.75, 11.5)
    textbox(s, "Posibles mejoras", 5.25, 1.35, 3.8, 0.3, 16, INK, True)
    bullets(s, [
        "Validar con meses futuros y monitorear drift.",
        "Elegir umbral por valor esperado: probabilidad x valor cliente x efectividad x costo.",
        "Agregar tickets, NPS, uso de servicio, promociones y cambios de plan.",
        "Calibrar probabilidades si el score se usa como ranking economico.",
        "Probar uplift modeling para contactar clientes persuadibles."
    ], 5.25, 1.82, 4.0, 2.75, 11.5)

    # 12 conclusiones
    s = new("9 - Conclusiones", "Un ranking accionable para priorizar retencion", 12)
    card(s, "Mejor modelo", "CatBoost", "mejor PR-AUC en CV", 0.7, 1.25, 2.0, 0.95, ORANGE)
    card(s, "PR-AUC test", num(metrics["test_pr_auc"]), "ranking", 3.0, 1.25, 1.8, 0.95, BLUE)
    card(s, "Recall", pct(metrics["test_recall"]), "churn detectado", 5.1, 1.25, 1.8, 0.95, GREEN)
    card(s, "Precision", pct(metrics["test_precision"]), "contactos utiles", 7.2, 1.25, 1.8, 0.95, RED)
    bullets(s, [
        "El modelo permite ordenar clientes por riesgo y priorizar una agenda de retencion.",
        "Los perfiles de mayor riesgo combinan contrato mensual, baja antiguedad y menor anclaje de servicios.",
        "Accion sugerida: contactar primero clientes de alto score con propuestas de permanencia, soporte o beneficios segun valor del cliente.",
        "Antes de produccion: definir costo de accion, capacidad del equipo y umbral por rentabilidad esperada."
    ], 0.9, 2.85, 8.2, 1.55, 13)

    prs.save(PPTX)

def main():
    df, metrics, cv, thresholds, preds, perm, shap, checks = load()
    make_assets(df, metrics, cv, thresholds, preds, perm, shap)
    build_deck(df, metrics, cv, preds, checks)
    print(PPTX)


if __name__ == "__main__":
    main()
