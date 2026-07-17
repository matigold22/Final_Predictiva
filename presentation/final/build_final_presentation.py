from pathlib import Path
import json
import textwrap

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "presentation" / "final"
ASSET_DIR = OUT_DIR / "assets"
ASSET_DIR.mkdir(parents=True, exist_ok=True)

DATA_PATH = ROOT / "data" / "raw" / "WA_Fn-UseC_-Telco-Customer-Churn.csv"
OUTPUTS = ROOT / "outputs"

PPTX_PATH = OUT_DIR / "Presentacion_Final_Churn_CatBoost.pptx"

COLORS = {
    "ink": RGBColor(28, 35, 49),
    "muted": RGBColor(92, 103, 121),
    "blue": RGBColor(44, 91, 190),
    "cyan": RGBColor(39, 166, 214),
    "green": RGBColor(69, 166, 111),
    "orange": RGBColor(238, 137, 55),
    "red": RGBColor(210, 78, 83),
    "paper": RGBColor(247, 249, 252),
    "white": RGBColor(255, 255, 255),
    "line": RGBColor(218, 225, 235),
}

MPL = {
    "ink": "#1c2331",
    "muted": "#5c6779",
    "blue": "#2c5bbe",
    "cyan": "#27a6d6",
    "green": "#45a66f",
    "orange": "#ee8937",
    "red": "#d24e53",
    "paper": "#f7f9fc",
    "line": "#dae1eb",
}


def pct(x, decimals=1):
    return f"{x * 100:.{decimals}f}%"


def fmt(x, decimals=3):
    return f"{x:.{decimals}f}"


def savefig(path):
    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close()


def load_data():
    df = pd.read_csv(DATA_PATH)
    df["TotalCharges"] = pd.to_numeric(df["TotalCharges"], errors="coerce").fillna(0)
    metrics = json.loads((OUTPUTS / "model_metrics.json").read_text(encoding="utf-8"))
    cv = pd.read_csv(OUTPUTS / "model_selection_cv_results.csv")
    thresholds = pd.read_csv(OUTPUTS / "threshold_search_results.csv").sort_values("threshold")
    preds = pd.read_csv(OUTPUTS / "test_predictions.csv")
    perm = pd.read_csv(OUTPUTS / "feature_importance_permutation_f2.csv")
    shap = pd.read_csv(OUTPUTS / "feature_importance_shap.csv")
    checks = json.loads((OUTPUTS / "leakage_overfit_checks.json").read_text(encoding="utf-8"))
    return df, metrics, cv, thresholds, preds, perm, shap, checks


def make_charts(df, metrics, cv, thresholds, preds, perm, shap, checks):
    dfp = df.copy()
    dfp["churn_num"] = (dfp["Churn"] == "Yes").astype(int)

    # 1. Churn and contract.
    fig, axes = plt.subplots(1, 2, figsize=(10, 3.6))
    churn = dfp["Churn"].value_counts(normalize=True).reindex(["No", "Yes"])
    axes[0].bar(["No churn", "Churn"], churn.values, color=[MPL["blue"], MPL["orange"]], width=0.55)
    axes[0].set_ylim(0, 1)
    axes[0].set_title("Distribucion del target", color=MPL["ink"], fontsize=12, weight="bold")
    axes[0].set_ylabel("Proporcion")
    for i, v in enumerate(churn.values):
        axes[0].text(i, v + 0.025, pct(v), ha="center", fontsize=11, weight="bold", color=MPL["ink"])

    contract = dfp.groupby("Contract")["churn_num"].mean().sort_values()
    axes[1].barh(contract.index, contract.values, color=MPL["green"])
    axes[1].set_xlim(0, max(contract.values) * 1.2)
    axes[1].set_title("Churn por tipo de contrato", color=MPL["ink"], fontsize=12, weight="bold")
    axes[1].set_xlabel("Proporcion churn")
    for i, v in enumerate(contract.values):
        axes[1].text(v + 0.01, i, pct(v), va="center", fontsize=10, weight="bold", color=MPL["ink"])
    for ax in axes:
        ax.spines[["top", "right"]].set_visible(False)
        ax.grid(axis="x", alpha=0.15)
    savefig(ASSET_DIR / "eda_target_contract.png")

    # 2. EDA segment signals.
    fig, axes = plt.subplots(1, 3, figsize=(12, 3.5))
    tenure_group = pd.cut(
        dfp["tenure"],
        bins=[-1, 6, 12, 24, 48, np.inf],
        labels=["0-6", "7-12", "13-24", "25-48", "49+"],
    )
    tg = dfp.assign(tenure_group=tenure_group).groupby("tenure_group", observed=False)["churn_num"].mean()
    axes[0].plot(tg.index.astype(str), tg.values, marker="o", color=MPL["blue"], linewidth=2.5)
    axes[0].set_title("Churn por antiguedad", fontsize=11, weight="bold")
    axes[0].set_ylim(0, max(tg.values) * 1.25)
    axes[0].set_ylabel("Proporcion churn")

    fiber = dfp.groupby("InternetService")["churn_num"].mean().sort_values()
    axes[1].barh(fiber.index, fiber.values, color=MPL["orange"])
    axes[1].set_title("Churn por internet", fontsize=11, weight="bold")

    support = dfp.groupby("TechSupport")["churn_num"].mean().sort_values()
    axes[2].barh(support.index, support.values, color=MPL["green"])
    axes[2].set_title("Churn por soporte tecnico", fontsize=11, weight="bold")

    for ax in axes:
        ax.spines[["top", "right"]].set_visible(False)
        ax.grid(axis="y" if ax == axes[0] else "x", alpha=0.15)
    savefig(ASSET_DIR / "eda_segments.png")

    # 3. Model selection PR-AUC.
    cvp = cv.sort_values("pr_auc_mean", ascending=True)
    fig, ax = plt.subplots(figsize=(8.4, 4.2))
    colors = [MPL["orange"] if m == metrics["selected_model"] else MPL["blue"] for m in cvp["model"]]
    ax.barh(cvp["model"], cvp["pr_auc_mean"], xerr=cvp["pr_auc_std"], color=colors, alpha=0.95)
    ax.axvline(dfp["churn_num"].mean(), color=MPL["red"], linestyle="--", linewidth=1.5, label="No skill = prevalencia")
    ax.set_title("Seleccion de modelos por PR-AUC (CV)", fontsize=13, weight="bold")
    ax.set_xlabel("PR-AUC promedio")
    ax.set_xlim(0.22, max(cvp["pr_auc_mean"]) + 0.04)
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(axis="x", alpha=0.18)
    ax.legend(frameon=False, loc="lower right")
    for i, (_, row) in enumerate(cvp.iterrows()):
        ax.text(row["pr_auc_mean"] + 0.004, i, f"{row['pr_auc_mean']:.3f}", va="center", fontsize=9)
    savefig(ASSET_DIR / "model_selection_prauc.png")

    # 4. Threshold tradeoff.
    fig, ax = plt.subplots(figsize=(8.6, 4.1))
    ax.plot(thresholds["threshold"], thresholds["recall"], label="Recall", color=MPL["blue"], linewidth=2.5)
    ax.plot(thresholds["threshold"], thresholds["precision"], label="Precision", color=MPL["orange"], linewidth=2.5)
    ax.plot(thresholds["threshold"], thresholds["f2"], label="F2", color=MPL["green"], linewidth=2.5)
    ax.axvline(metrics["decision_threshold"], color=MPL["red"], linestyle="--", linewidth=1.5)
    ax.text(metrics["decision_threshold"] + 0.012, 0.08, f"umbral {metrics['decision_threshold']:.2f}", color=MPL["red"], fontsize=10, weight="bold")
    ax.set_ylim(0, 1.02)
    ax.set_xlim(0.1, 0.9)
    ax.set_title("Trade-off del umbral: recall vs precision", fontsize=13, weight="bold")
    ax.set_xlabel("Umbral de probabilidad")
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(alpha=0.18)
    ax.legend(frameon=False, loc="upper right")
    savefig(ASSET_DIR / "threshold_tradeoff.png")

    # 5. Confusion matrix.
    cm = pd.crosstab(preds["y_true_label"], preds["churn_prediction_label"]).reindex(index=["No", "Yes"], columns=["No", "Yes"], fill_value=0)
    fig, ax = plt.subplots(figsize=(4.8, 4.3))
    im = ax.imshow(cm.values, cmap="Blues")
    ax.set_xticks([0, 1], ["Pred No", "Pred Churn"])
    ax.set_yticks([0, 1], ["Real No", "Real Churn"])
    ax.set_title("Matriz de confusion en test", fontsize=12, weight="bold")
    for i in range(2):
        for j in range(2):
            ax.text(j, i, str(cm.values[i, j]), ha="center", va="center", fontsize=15, weight="bold", color=MPL["ink"])
    fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    savefig(ASSET_DIR / "confusion_matrix.png")

    # 6. Lift / capture curve.
    ranked = preds.sort_values("churn_probability", ascending=False).copy()
    ranked["cum_churn"] = ranked["y_true"].cumsum()
    total_churn = ranked["y_true"].sum()
    ranked["portfolio_share"] = np.arange(1, len(ranked) + 1) / len(ranked)
    ranked["captured_churn_share"] = ranked["cum_churn"] / total_churn
    fig, ax = plt.subplots(figsize=(7.2, 4.1))
    ax.plot(ranked["portfolio_share"], ranked["captured_churn_share"], color=MPL["blue"], linewidth=2.8, label="Modelo")
    ax.plot([0, 1], [0, 1], color=MPL["muted"], linestyle="--", label="Azar")
    for cut in [0.2, 0.4, 0.6]:
        captured = ranked.loc[ranked["portfolio_share"] <= cut, "y_true"].sum() / total_churn
        ax.scatter([cut], [captured], color=MPL["orange"], zorder=3)
        ax.text(cut + 0.02, captured, f"Top {int(cut*100)}% -> {pct(captured, 0)}", fontsize=9, weight="bold")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_title("Curva de captura de churners", fontsize=13, weight="bold")
    ax.set_xlabel("Porcentaje de cartera accionada")
    ax.set_ylabel("Churners capturados")
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(alpha=0.18)
    ax.legend(frameon=False, loc="lower right")
    savefig(ASSET_DIR / "lift_capture_curve.png")

    # 7. Feature importance permutation.
    top = perm.head(10).sort_values("importance_mean")
    fig, ax = plt.subplots(figsize=(7.4, 4.6))
    ax.barh(top["feature"], top["importance_mean"], xerr=top["importance_std"], color=MPL["green"])
    ax.set_title("Permutation importance alineada a F2", fontsize=13, weight="bold")
    ax.set_xlabel("Caida de F2 al desordenar la variable")
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(axis="x", alpha=0.18)
    savefig(ASSET_DIR / "feature_importance_perm_clean.png")

    # 8. SHAP.
    top = shap.head(10).sort_values("mean_abs_shap")
    fig, ax = plt.subplots(figsize=(7.4, 4.6))
    ax.barh(top["feature"], top["mean_abs_shap"], color=MPL["orange"])
    ax.set_title("SHAP global: impacto medio absoluto", fontsize=13, weight="bold")
    ax.set_xlabel("mean(|SHAP|)")
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(axis="x", alpha=0.18)
    savefig(ASSET_DIR / "feature_importance_shap_clean.png")


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_bg(slide, color=COLORS["paper"]):
    shape = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(5.625))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_footer(slide, n):
    tx = slide.shapes.add_textbox(Inches(0.45), Inches(5.25), Inches(9.1), Inches(0.2))
    p = tx.text_frame.paragraphs[0]
    p.text = f"Analisis Predictivo · Churn Telco · ITBA 2026 · {n:02d}"
    p.font.size = Pt(7.5)
    p.font.color.rgb = COLORS["muted"]


def add_title(slide, section, title, subtitle=None):
    sec = slide.shapes.add_textbox(Inches(0.45), Inches(0.28), Inches(2.2), Inches(0.25))
    p = sec.text_frame.paragraphs[0]
    p.text = section.upper()
    p.font.size = Pt(8.5)
    p.font.bold = True
    p.font.color.rgb = COLORS["blue"]

    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.55), Inches(8.9), Inches(0.65))
    tf = box.text_frame
    tf.margin_left = 0
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = COLORS["ink"]
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLORS["muted"]


def add_text(slide, text, x, y, w, h, size=12, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS["ink"]
    if align:
        p.alignment = align
    return box


def add_bullets(slide, items, x, y, w, h, size=12, color=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color or COLORS["ink"]
        p.space_after = Pt(5)
    return box


def add_card(slide, x, y, w, h, title, value=None, body=None, accent=COLORS["blue"]):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(1)

    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    add_text(slide, title, x + 0.18, y + 0.12, w - 0.28, 0.22, size=8.5, bold=True, color=COLORS["muted"])
    if value is not None:
        add_text(slide, value, x + 0.18, y + 0.38, w - 0.28, 0.45, size=22, bold=True, color=COLORS["ink"])
    if body:
        add_text(slide, body, x + 0.18, y + (0.85 if value else 0.38), w - 0.28, h - 0.9, size=9.5, color=COLORS["muted"])


def add_image(slide, path, x, y, w, h=None):
    if h is None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_table_like(slide, rows, x, y, w, h, col_widths=None, font_size=9):
    n_rows = len(rows)
    n_cols = len(rows[0])
    table = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["blue"] if r == 0 else COLORS["white"]
            cell.text_frame.margin_left = Inches(0.04)
            cell.text_frame.margin_right = Inches(0.04)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(font_size)
            p.font.bold = r == 0
            p.font.color.rgb = COLORS["white"] if r == 0 else COLORS["ink"]
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
    return table


def build_deck(df, metrics, cv, thresholds, preds, perm, shap, checks):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    slide_no = 0

    def new(section, title, subtitle=None):
        nonlocal slide_no
        slide_no += 1
        slide = blank_slide(prs)
        add_bg(slide)
        add_title(slide, section, title, subtitle)
        add_footer(slide, slide_no)
        return slide

    # 1
    slide_no += 1
    slide = blank_slide(prs)
    add_bg(slide, COLORS["ink"])
    add_text(slide, "ANALISIS PREDICTIVO - EXAMEN FINAL", 0.55, 0.48, 5.5, 0.3, 10, True, COLORS["cyan"])
    add_text(slide, "Prediccion de churn en una telco", 0.55, 1.05, 6.7, 0.9, 30, True, COLORS["white"])
    add_text(slide, "Un modelo para priorizar clientes con riesgo de baja antes de que se vayan.", 0.58, 1.95, 5.9, 0.55, 14, False, RGBColor(202, 213, 229))
    add_card(slide, 0.65, 3.35, 2.0, 0.95, "CHURN", pct((df["Churn"] == "Yes").mean()), "clase positiva", COLORS["orange"])
    add_card(slide, 2.9, 3.35, 2.0, 0.95, "DATASET", f"{len(df):,}".replace(",", "."), "clientes", COLORS["cyan"])
    add_card(slide, 5.15, 3.35, 2.0, 0.95, "MODELO FINAL", "CatBoost", "PR-AUC + F2", COLORS["green"])
    add_text(slide, "Felipe Tamaki | ITBA · 2026 · IBM Telco Customer Churn", 0.6, 5.1, 6.5, 0.25, 8.5, False, RGBColor(202, 213, 229))

    # 2
    slide = new("1 · Caso de negocio", "Retener cuesta menos que adquirir: hay que llegar antes")
    add_bullets(slide, [
        "Problema: identificar clientes con mayor probabilidad de baja.",
        "Decision: priorizar acciones de retencion sobre una cartera limitada.",
        "Valor: transformar datos contractuales, servicios y facturacion en un ranking accionable.",
        "Trade-off: preferimos alto recall, pero controlando precision para no disparar campañas enormes."
    ], 0.55, 1.55, 5.3, 2.3, 14)
    add_card(slide, 6.1, 1.45, 1.65, 0.9, "TASA CHURN", pct((df["Churn"] == "Yes").mean()), None, COLORS["orange"])
    add_card(slide, 7.95, 1.45, 1.45, 0.9, "CLIENTES", f"{len(df):,}".replace(",", "."), None, COLORS["blue"])
    add_card(slide, 6.1, 2.65, 3.3, 1.1, "OUTPUT DEL MODELO", "score de riesgo", "probabilidad de churn por cliente", COLORS["green"])
    add_card(slide, 6.1, 4.0, 3.3, 0.75, "USO", None, "ordenar cartera y decidir a quien contactar", COLORS["cyan"])

    # 3
    slide = new("1 · Caso de negocio", "Flujo operativo: del dato mensual a la agenda de retencion")
    steps = [
        ("1", "Base activa", "contrato, antiguedad, servicios, facturacion"),
        ("2", "Modelo", "score de churn para cada cliente"),
        ("3", "Umbral", "prioriza recall con F2"),
        ("4", "Accion", "retencion, oferta o contacto"),
    ]
    x = 0.6
    for n, title, body in steps:
        add_card(slide, x, 1.65, 1.85, 1.7, title, n, body, COLORS["blue"] if n != "4" else COLORS["green"])
        if n != "4":
            add_text(slide, "→", x + 1.92, 2.23, 0.35, 0.3, 24, True, COLORS["muted"])
        x += 2.25
    add_text(slide, "La prediccion no reemplaza al negocio: entrega un ranking. El corte final deberia depender del costo de contactar, el valor del cliente y la efectividad esperada de la accion.", 0.85, 4.2, 8.25, 0.55, 13, False, COLORS["muted"], PP_ALIGN.CENTER)

    # 4
    slide = new("2 · Dataset", "IBM Telco Customer Churn: 7.043 clientes y 21 variables")
    add_card(slide, 0.55, 1.42, 2.0, 1.25, "DEMOGRAFICAS", "4", "genero, senior, pareja, dependientes", COLORS["blue"])
    add_card(slide, 2.85, 1.42, 2.0, 1.25, "SERVICIOS", "9", "internet, seguridad, soporte, streaming", COLORS["green"])
    add_card(slide, 5.15, 1.42, 2.0, 1.25, "CONTRATO", "4", "tenure, contract, payment, paperless", COLORS["orange"])
    add_card(slide, 7.45, 1.42, 2.0, 1.25, "FACTURACION", "2", "MonthlyCharges, TotalCharges", COLORS["cyan"])
    add_bullets(slide, [
        "Unidad de analisis: un cliente.",
        "Target: Churn = Yes/No.",
        "Limpieza relevante: TotalCharges venia como texto; 11 clientes con tenure = 0 se imputan en 0.",
        "Limitacion: foto de un mes, sin historia temporal de uso, reclamos ni satisfaccion."
    ], 0.75, 3.08, 8.5, 1.45, 12.5)

    # 5
    slide = new("3 · Analisis exploratorio", "El churn no aparece al azar: contrato y antiguedad concentran señal")
    add_image(slide, ASSET_DIR / "eda_target_contract.png", 0.45, 1.35, 4.7)
    add_image(slide, ASSET_DIR / "eda_segments.png", 5.05, 1.35, 4.55)
    add_text(slide, "Lectura: contrato mensual, baja antiguedad, fibra optica y ausencia de soporte son señales consistentes. No son causalidad: son patrones predictivos.", 0.75, 4.85, 8.4, 0.35, 10.5, False, COLORS["muted"], PP_ALIGN.CENTER)

    # 6
    slide = new("4 · Evaluacion", "Con clases desbalanceadas, accuracy no alcanza")
    add_card(slide, 0.65, 1.35, 2.05, 1.15, "NO SKILL PR-AUC", fmt((df["Churn"] == "Yes").mean()), "igual a la prevalencia", COLORS["red"])
    add_card(slide, 2.95, 1.35, 2.05, 1.15, "TRAIN", "5.634", "80%, estratificado", COLORS["blue"])
    add_card(slide, 5.25, 1.35, 2.05, 1.15, "TEST", "1.409", "20%, cerrado", COLORS["cyan"])
    add_card(slide, 7.55, 1.35, 1.8, 1.15, "SEED", "42", "reproducible", COLORS["green"])
    add_bullets(slide, [
        "PR-AUC para seleccionar modelos: evalua ranking de churners y se enfoca en la clase positiva.",
        "F2 para elegir umbral: prioriza recall sin ignorar precision.",
        "Validacion cruzada estratificada de 3 folds sobre train para todos los modelos.",
        "El test se usa una sola vez para evaluacion final."
    ], 0.75, 3.0, 8.6, 1.35, 12.5)

    # 7
    slide = new("5 · Modelo baseline", "Baseline: el piso minimo y una referencia interpretable")
    base = cv.set_index("model")
    add_card(slide, 0.65, 1.45, 2.15, 1.2, "DUMMY", fmt(base.loc["dummy_most_frequent", "pr_auc_mean"]), "PR-AUC = prevalencia", COLORS["red"])
    add_card(slide, 3.1, 1.45, 2.15, 1.2, "LOGISTIC", fmt(base.loc["logistic_regression", "pr_auc_mean"]), "PR-AUC CV", COLORS["blue"])
    add_card(slide, 5.55, 1.45, 2.15, 1.2, "LOGISTIC F2", fmt(base.loc["logistic_regression", "f2_mean"]), "muy competitivo", COLORS["green"])
    add_bullets(slide, [
        "La regresion logistica queda muy cerca del ganador: buena señal de que el problema tiene estructura clara.",
        "Se mantiene como benchmark interpretable: si un modelo complejo gana por poco, debe justificar su complejidad.",
        "El baseline tonto confirma por que accuracy no sirve: puede acertar muchos No Churn y no detectar ningun churner."
    ], 0.9, 3.08, 8.1, 1.35, 12.5)

    # 8
    slide = new("6 · Seleccion de modelos", "CatBoost gana por PR-AUC, pero la carrera fue pareja")
    add_image(slide, ASSET_DIR / "model_selection_prauc.png", 0.55, 1.35, 5.1)
    top3 = cv.sort_values("pr_auc_mean", ascending=False).head(3)
    rows = [["Modelo", "PR-AUC", "F2", "Recall", "Precision"]]
    for _, r in top3.iterrows():
        rows.append([r["model"], fmt(r["pr_auc_mean"]), fmt(r["f2_mean"]), fmt(r["recall_mean"]), fmt(r["precision_mean"])])
    add_table_like(slide, rows, 5.85, 1.55, 3.75, 1.45, [1.4, 0.58, 0.48, 0.58, 0.7], 8)
    add_text(slide, "Top 3 separados por menos de 0,005 PR-AUC: CatBoost gana, pero Logistic sigue siendo un respaldo simple e interpretable.", 5.9, 3.35, 3.55, 0.85, 11.5, False, COLORS["muted"])

    # 9
    slide = new("6 · Seleccion de modelos", "Tuning acotado: todos los modelos con CV y misma metrica")
    params = metrics["best_params"]
    add_card(slide, 0.65, 1.45, 2.0, 1.05, "MODELO", "CatBoost", "mejor PR-AUC CV", COLORS["orange"])
    add_card(slide, 2.95, 1.45, 1.35, 1.05, "DEPTH", str(params["model__depth"]), None, COLORS["blue"])
    add_card(slide, 4.55, 1.45, 1.65, 1.05, "ITER.", str(params["model__iterations"]), None, COLORS["green"])
    add_card(slide, 6.45, 1.45, 1.45, 1.05, "LR", str(params["model__learning_rate"]), None, COLORS["cyan"])
    add_card(slide, 8.15, 1.45, 1.25, 1.05, "L2", str(params["model__l2_leaf_reg"]), None, COLORS["red"])
    add_bullets(slide, [
        "Se probaron modelos lineales, arboles, Random Forest, Gradient Boosting, LightGBM y CatBoost.",
        "La busqueda fue acotada para mantener reproducibilidad y evitar tuning excesivo sobre una base chica.",
        "LightGBM fue probado, pero quedo por debajo en PR-AUC CV.",
        "CatBoost combina buen ranking con brecha train-test baja a moderada."
    ], 0.85, 3.0, 8.25, 1.45, 12.5)

    # 10
    slide = new("7 · Modelo final", "CatBoost en test: ranking fuerte y foco en recall")
    add_card(slide, 0.65, 1.35, 1.55, 1.0, "PR-AUC", fmt(metrics["test_pr_auc"]), "test", COLORS["orange"])
    add_card(slide, 2.45, 1.35, 1.55, 1.0, "ROC-AUC", fmt(metrics["test_roc_auc"]), "test", COLORS["blue"])
    add_card(slide, 4.25, 1.35, 1.55, 1.0, "Recall", pct(metrics["test_recall"]), "churn detectado", COLORS["green"])
    add_card(slide, 6.05, 1.35, 1.55, 1.0, "Precision", pct(metrics["test_precision"]), "contactos utiles", COLORS["red"])
    add_card(slide, 7.85, 1.35, 1.55, 1.0, "F2", fmt(metrics["test_f2"]), "umbral 0.22", COLORS["cyan"])
    add_bullets(slide, [
        "El modelo final no se eligio por accuracy, sino por PR-AUC.",
        "El umbral 0.22 maximiza F2 en validacion interna: prioriza detectar churners.",
        "La precision queda baja: es un trade-off operativo consciente, no una metrica a esconder.",
        f"El modelo marca como riesgo al {pct(metrics['positive_rate_predicted_test'])} del test; requiere una estrategia de accion por capacidad/costos."
    ], 0.85, 3.0, 8.25, 1.55, 12.2)

    # 11
    slide = new("7 · Modelo final", "El umbral define la politica de accion")
    add_image(slide, ASSET_DIR / "threshold_tradeoff.png", 0.55, 1.35, 5.25)
    add_text(slide, "Umbral elegido: 0.22", 6.1, 1.45, 2.9, 0.35, 20, True, COLORS["ink"])
    add_bullets(slide, [
        "Umbral bajo => recall alto y mas falsos positivos.",
        "Con costos reales, el corte deberia elegirse por valor esperado.",
        "PR-AUC elige el mejor ranking; F2 transforma ese ranking en una decision."
    ], 6.15, 2.1, 3.1, 1.45, 12)
    add_card(slide, 6.15, 3.9, 1.45, 0.75, "Recall", pct(metrics["test_recall"]), None, COLORS["green"])
    add_card(slide, 7.85, 3.9, 1.45, 0.75, "Precision", pct(metrics["test_precision"]), None, COLORS["red"])

    # 12
    slide = new("7 · Modelo final", "Impacto operativo: muchos churners capturados, cartera amplia")
    add_image(slide, ASSET_DIR / "confusion_matrix.png", 0.65, 1.35, 3.1)
    add_image(slide, ASSET_DIR / "lift_capture_curve.png", 4.15, 1.35, 5.05)
    add_text(slide, "La matriz muestra el costo del umbral agresivo: muchos falsos positivos. La curva de captura permite definir cortes alternativos si el equipo comercial tiene capacidad limitada.", 0.85, 4.78, 8.35, 0.35, 10.5, False, COLORS["muted"], PP_ALIGN.CENTER)

    # 13
    slide = new("7 · Interpretacion", "Variables importantes: contrato, antiguedad y servicios de anclaje")
    add_image(slide, ASSET_DIR / "feature_importance_perm_clean.png", 0.55, 1.3, 4.7)
    add_image(slide, ASSET_DIR / "feature_importance_shap_clean.png", 5.15, 1.3, 4.35)
    add_text(slide, "Permutation importance esta alineada con F2; SHAP ayuda a leer impacto global. Ambas son señales predictivas, no causalidad.", 0.75, 4.86, 8.5, 0.35, 10.5, False, COLORS["muted"], PP_ALIGN.CENTER)

    # 14
    slide = new("7 · Interpretacion", "Lectura de negocio: perfiles a mirar primero")
    add_bullets(slide, [
        "Contrato mensual: la señal dominante; menos compromiso contractual.",
        "Baja antiguedad: los clientes nuevos concentran mayor riesgo.",
        "Fibra optica: aparece asociada a churn; puede capturar precio, experiencia o expectativas.",
        "Sin seguridad/soporte tecnico: menor anclaje al servicio.",
        "Cargos acumulados y gap de gasto: resumen historia y relacion con facturacion actual."
    ], 0.75, 1.45, 4.7, 2.65, 14)
    add_card(slide, 5.85, 1.55, 3.2, 0.9, "CUIDADO", None, "No interpretamos como causalidad: son variables que ayudan a predecir.", COLORS["red"])
    add_card(slide, 5.85, 2.75, 3.2, 0.9, "ACCION", None, "Usar score + criterio comercial para priorizar retencion.", COLORS["green"])
    add_card(slide, 5.85, 3.95, 3.2, 0.9, "VALIDACION", None, "Monitorear precision/recall y recalibrar con datos futuros.", COLORS["blue"])

    # 15
    slide = new("8 · Limitaciones y mejoras", "El techo no esta solo en el algoritmo")
    add_card(slide, 0.65, 1.35, 2.55, 1.0, "SIN TIEMPO", None, "no hay corte pasado-futuro", COLORS["red"])
    add_card(slide, 3.55, 1.35, 2.55, 1.0, "VARIABLES REDUNDANTES", None, "feature engineering + originales", COLORS["orange"])
    add_card(slide, 6.45, 1.35, 2.55, 1.0, "UMBRAL", None, "requiere costos reales", COLORS["blue"])
    add_bullets(slide, [
        "Validar con meses posteriores para medir estabilidad temporal.",
        "Incorporar reclamos, uso real, tickets de soporte, promociones y cambios de plan.",
        "Elegir umbral por valor esperado: probabilidad x valor cliente x efectividad - costo.",
        "Probar uplift modeling para contactar persuadibles, no solo churners probables.",
        "Reducir redundancia si se prioriza interpretabilidad."
    ], 0.85, 2.85, 8.4, 1.65, 12.5)

    # 16
    slide = new("9 · Conclusiones", "Un problema de negocio resuelto con metodo")
    add_card(slide, 0.65, 1.35, 2.55, 1.25, "SENAL CLARA", None, "contrato, tenure, internet y soporte aparecen en EDA y modelo", COLORS["green"])
    add_card(slide, 3.75, 1.35, 2.55, 1.25, "METODO HONESTO", None, "PR-AUC para ranking, F2 para umbral, test usado al final", COLORS["blue"])
    add_card(slide, 6.85, 1.35, 2.55, 1.25, "MODELO FINAL", None, "CatBoost con PR-AUC 0.667 y recall 95.5%", COLORS["orange"])
    add_bullets(slide, [
        "El modelo sirve para priorizar retencion, no para explicar causalidad.",
        "El punto operativo actual es agresivo: captura churners pero genera falsos positivos.",
        "La siguiente mejora no es otro algoritmo: es validar con datos futuros y definir costos de negocio."
    ], 1.0, 3.15, 8.0, 1.3, 14)

    prs.save(PPTX_PATH)


def main():
    df, metrics, cv, thresholds, preds, perm, shap, checks = load_data()
    make_charts(df, metrics, cv, thresholds, preds, perm, shap, checks)
    build_deck(df, metrics, cv, thresholds, preds, perm, shap, checks)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
