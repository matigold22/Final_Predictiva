from pathlib import Path
import json

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "presentation" / "final"
ASSETS = OUT / "assets_redesign"
ASSETS.mkdir(parents=True, exist_ok=True)
PPTX = OUT / "Presentacion_Final_Churn_Redesign.pptx"

DATA = ROOT / "data" / "raw" / "WA_Fn-UseC_-Telco-Customer-Churn.csv"
OUTPUTS = ROOT / "outputs"

INK = RGBColor(22, 28, 39)
MUTED = RGBColor(94, 104, 119)
PAPER = RGBColor(249, 250, 252)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(220, 226, 235)
BLUE = RGBColor(35, 86, 194)
TEAL = RGBColor(34, 151, 153)
ORANGE = RGBColor(231, 126, 55)
RED = RGBColor(200, 71, 82)
GREEN = RGBColor(60, 154, 96)
VIOLET = RGBColor(104, 86, 204)

M = {
    "ink": "#161c27",
    "muted": "#5e6877",
    "paper": "#f9fafc",
    "line": "#dce2eb",
    "blue": "#2356c2",
    "teal": "#229799",
    "orange": "#e77e37",
    "red": "#c84752",
    "green": "#3c9a60",
    "violet": "#6856cc",
}


def pct(x, d=1):
    return f"{100*x:.{d}f}%"


def num(x, d=3):
    return f"{x:.{d}f}"


def save(path):
    plt.tight_layout()
    plt.savefig(path, dpi=190, bbox_inches="tight", facecolor="white")
    plt.close()


def load():
    df = pd.read_csv(DATA)
    df["TotalCharges"] = pd.to_numeric(df["TotalCharges"], errors="coerce").fillna(0)
    df["churn_num"] = (df["Churn"] == "Yes").astype(int)
    metrics = json.loads((OUTPUTS / "model_metrics.json").read_text(encoding="utf-8"))
    cv = pd.read_csv(OUTPUTS / "model_selection_cv_results.csv")
    thresholds = pd.read_csv(OUTPUTS / "threshold_search_results.csv").sort_values("threshold")
    preds = pd.read_csv(OUTPUTS / "test_predictions.csv")
    perm = pd.read_csv(OUTPUTS / "feature_importance_permutation_f2.csv")
    shap = pd.read_csv(OUTPUTS / "feature_importance_shap.csv")
    checks = json.loads((OUTPUTS / "leakage_overfit_checks.json").read_text(encoding="utf-8"))
    return df, metrics, cv, thresholds, preds, perm, shap, checks


def style_ax(ax):
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(alpha=0.16)
    ax.tick_params(colors=M["muted"], labelsize=9)
    ax.title.set_color(M["ink"])
    ax.xaxis.label.set_color(M["muted"])
    ax.yaxis.label.set_color(M["muted"])


def make_assets(df, metrics, cv, thresholds, preds, perm, shap):
    # EDA: churn rate vs multiple variables/categories.
    variables = {
        "Contract": "Contrato",
        "InternetService": "Internet",
        "PaymentMethod": "Pago",
        "OnlineSecurity": "Seguridad online",
        "TechSupport": "Soporte tecnico",
        "PaperlessBilling": "Factura digital",
        "Dependents": "Dependientes",
        "Partner": "Pareja",
    }
    rows = []
    for col, label in variables.items():
        rates = df.groupby(col)["churn_num"].mean().sort_values()
        for cat, rate in rates.items():
            rows.append({"variable": label, "category": str(cat), "rate": rate})
    plot = pd.DataFrame(rows)
    order = (
        plot.groupby("variable")["rate"].agg(lambda s: s.max() - s.min())
        .sort_values(ascending=True)
        .index.tolist()
    )
    y_map = {v: i for i, v in enumerate(order)}
    fig, ax = plt.subplots(figsize=(9.2, 5.0))
    for variable in order:
        sub = plot[plot["variable"] == variable]
        y = y_map[variable]
        ax.plot([sub["rate"].min(), sub["rate"].max()], [y, y], color=M["line"], linewidth=4, zorder=1)
        for _, r in sub.iterrows():
            color = M["orange"] if r["rate"] >= df["churn_num"].mean() else M["blue"]
            ax.scatter(r["rate"], y, s=95, color=color, edgecolor="white", linewidth=1.2, zorder=3)
            ax.text(r["rate"] + 0.009, y + 0.06, r["category"][:18], fontsize=8, color=M["muted"])
    ax.axvline(df["churn_num"].mean(), color=M["red"], linestyle="--", linewidth=1.5)
    ax.text(df["churn_num"].mean() + 0.01, len(order) - 0.45, f"media {pct(df['churn_num'].mean())}", color=M["red"], fontsize=9, weight="bold")
    ax.set_yticks(range(len(order)), order)
    ax.set_xlim(0, 0.5)
    ax.set_xlabel("Tasa de churn por categoria")
    ax.set_title("Tasa de churn cruzada por variables relevantes", fontsize=14, weight="bold")
    style_ax(ax)
    save(ASSETS / "eda_churn_rate_variables.png")

    # Model PR-AUC with 10 tests.
    cvp = cv.sort_values("pr_auc_mean", ascending=True).copy()
    labels = {
        "dummy_most_frequent": "Dummy",
        "gaussian_nb": "Gaussian NB",
        "decision_tree": "Decision Tree",
        "ada_boost": "AdaBoost",
        "extra_trees": "Extra Trees",
        "lightgbm": "LightGBM",
        "random_forest": "Random Forest",
        "logistic_regression": "Logistic Reg.",
        "gradient_boosting": "Gradient Boosting",
        "catboost": "CatBoost",
    }
    cvp["label"] = cvp["model"].map(labels).fillna(cvp["model"])
    fig, ax = plt.subplots(figsize=(8.6, 5.0))
    colors = [M["orange"] if m == metrics["selected_model"] else M["blue"] for m in cvp["model"]]
    ax.barh(cvp["label"], cvp["pr_auc_mean"], xerr=cvp["pr_auc_std"], color=colors, height=0.62)
    ax.axvline(df["churn_num"].mean(), color=M["red"], linestyle="--", linewidth=1.2)
    ax.text(df["churn_num"].mean() + 0.006, -0.45, "no skill", color=M["red"], fontsize=8, weight="bold")
    for i, r in enumerate(cvp.itertuples()):
        ax.text(r.pr_auc_mean + 0.004, i, f"{r.pr_auc_mean:.3f}", va="center", fontsize=8.5, color=M["ink"])
    ax.set_xlim(0.23, max(cvp["pr_auc_mean"]) + 0.045)
    ax.set_xlabel("PR-AUC promedio en validacion cruzada")
    ax.set_title("10 pruebas de modelos: comparacion por PR-AUC", fontsize=14, weight="bold")
    style_ax(ax)
    save(ASSETS / "modelos_10_prauc.png")

    # Compact EDA target + tenure.
    fig, axes = plt.subplots(1, 2, figsize=(9, 3.7))
    churn = df["Churn"].value_counts(normalize=True).reindex(["No", "Yes"])
    axes[0].bar(["No churn", "Churn"], churn.values, color=[M["blue"], M["orange"]], width=0.55)
    axes[0].set_ylim(0, 1)
    axes[0].set_title("Distribucion del target", fontsize=12, weight="bold")
    for i, v in enumerate(churn.values):
        axes[0].text(i, v + 0.025, pct(v), ha="center", fontsize=11, weight="bold", color=M["ink"])
    tg = df.assign(tenure_group=pd.cut(df["tenure"], [-1, 6, 12, 24, 48, np.inf], labels=["0-6", "7-12", "13-24", "25-48", "49+"]))
    rates = tg.groupby("tenure_group", observed=False)["churn_num"].mean()
    axes[1].plot(rates.index.astype(str), rates.values, marker="o", color=M["teal"], linewidth=2.7)
    axes[1].set_title("Churn por antiguedad", fontsize=12, weight="bold")
    axes[1].set_ylabel("Tasa churn")
    for ax in axes:
        style_ax(ax)
    save(ASSETS / "eda_target_tenure.png")

    # Threshold.
    fig, ax = plt.subplots(figsize=(8.2, 4.3))
    ax.plot(thresholds["threshold"], thresholds["recall"], color=M["blue"], linewidth=2.5, label="Recall")
    ax.plot(thresholds["threshold"], thresholds["precision"], color=M["orange"], linewidth=2.5, label="Precision")
    ax.plot(thresholds["threshold"], thresholds["f2"], color=M["green"], linewidth=2.5, label="F2")
    ax.axvline(metrics["decision_threshold"], color=M["red"], linestyle="--")
    ax.set_xlim(0.1, 0.9)
    ax.set_ylim(0, 1.02)
    ax.set_xlabel("Umbral")
    ax.set_title("Umbral operativo: trade-off precision / recall", fontsize=14, weight="bold")
    ax.legend(frameon=False, loc="upper right")
    style_ax(ax)
    save(ASSETS / "threshold_precision_recall.png")

    # Feature importance.
    top = perm.head(10).sort_values("importance_mean")
    fig, ax = plt.subplots(figsize=(7.7, 4.7))
    ax.barh(top["feature"], top["importance_mean"], xerr=top["importance_std"], color=M["teal"], height=0.6)
    ax.set_xlabel("Caida de F2 al permutar")
    ax.set_title("Feature importance por permutation importance", fontsize=14, weight="bold")
    style_ax(ax)
    save(ASSETS / "importance_permutation.png")

    top = shap.head(10).sort_values("mean_abs_shap")
    fig, ax = plt.subplots(figsize=(7.7, 4.7))
    ax.barh(top["feature"], top["mean_abs_shap"], color=M["violet"], height=0.6)
    ax.set_xlabel("mean(|SHAP|)")
    ax.set_title("SHAP global del modelo final", fontsize=14, weight="bold")
    style_ax(ax)
    save(ASSETS / "importance_shap.png")

    # Capture curve.
    ranked = preds.sort_values("churn_probability", ascending=False).copy()
    ranked["portfolio_share"] = np.arange(1, len(ranked) + 1) / len(ranked)
    ranked["captured_churn"] = ranked["y_true"].cumsum() / ranked["y_true"].sum()
    fig, ax = plt.subplots(figsize=(7.8, 4.3))
    ax.plot(ranked["portfolio_share"], ranked["captured_churn"], color=M["blue"], linewidth=2.8)
    ax.plot([0, 1], [0, 1], color=M["muted"], linestyle="--")
    for cut in [0.2, 0.4, 0.6]:
        captured = ranked.loc[ranked["portfolio_share"] <= cut, "y_true"].sum() / ranked["y_true"].sum()
        ax.scatter([cut], [captured], color=M["orange"], s=70)
        ax.text(cut + 0.02, captured, f"Top {int(cut*100)}%: {pct(captured,0)}", fontsize=9, weight="bold")
    ax.set_xlabel("Porcentaje de cartera accionada")
    ax.set_ylabel("Churners capturados")
    ax.set_title("Curva de captura de churners", fontsize=14, weight="bold")
    style_ax(ax)
    save(ASSETS / "capture_curve.png")


def add_bg(slide):
    bg = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(5.625))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    band = slide.shapes.add_shape(1, 0, 0, Inches(0.18), Inches(5.625))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()


def textbox(slide, text, x, y, w, h, size=12, color=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def title(slide, section, main):
    textbox(slide, section.upper(), 0.48, 0.24, 3.1, 0.28, 8.5, BLUE, True)
    textbox(slide, main, 0.48, 0.52, 8.9, 0.65, 23, INK, True)


def bullets(slide, items, x, y, w, h, size=12):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(5)
    return box


def card(slide, label, value, note, x, y, w, h, color=BLUE):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    stripe = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.08))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()
    textbox(slide, label, x + 0.15, y + 0.16, w - 0.3, 0.22, 8.2, MUTED, True)
    if value:
        textbox(slide, value, x + 0.15, y + 0.42, w - 0.3, 0.42, 20, INK, True)
    if note:
        textbox(slide, note, x + 0.15, y + (0.86 if value else 0.42), w - 0.3, h - 0.9, 9.2, MUTED)


def image(slide, path, x, y, w, h=None):
    if h:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))


def footer(slide, n):
    textbox(slide, f"Churn Telco · Examen Final · {n:02d}", 0.48, 5.28, 4.5, 0.18, 7.5, MUTED)


def build_deck(df, metrics, cv, preds, checks):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    def new(section, main, n):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(s)
        title(s, section, main)
        footer(s, n)
        return s

    # 1 portada
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, Inches(10), Inches(5.625))
    bg.fill.solid(); bg.fill.fore_color.rgb = INK; bg.line.fill.background()
    textbox(s, "ANALISIS PREDICTIVO · EXAMEN FINAL", 0.6, 0.55, 5.5, 0.3, 9.5, RGBColor(146, 196, 255), True)
    textbox(s, "Prediccion de churn en clientes Telco", 0.6, 1.15, 6.8, 0.95, 31, WHITE, True)
    textbox(s, "Modelo final CatBoost · seleccion por PR-AUC · umbral optimizado por F2", 0.64, 2.1, 6.8, 0.35, 13, RGBColor(205, 215, 230))
    card(s, "CHURN", pct(df["churn_num"].mean()), "tasa observada", 0.75, 3.55, 1.7, 0.95, ORANGE)
    card(s, "CLIENTES", f"{len(df):,}".replace(",", "."), "dataset IBM", 2.75, 3.55, 1.7, 0.95, TEAL)
    card(s, "PR-AUC TEST", num(metrics["test_pr_auc"]), "modelo final", 4.75, 3.55, 1.7, 0.95, GREEN)
    card(s, "RECALL", pct(metrics["test_recall"]), "umbral 0.31", 6.75, 3.55, 1.7, 0.95, BLUE)

    # 2 caso
    s = new("1 · Caso de negocio", "Priorizar retencion antes de que ocurra la baja", 2)
    bullets(s, [
        "Objetivo: estimar riesgo de churn por cliente activo.",
        "Decision: ordenar la cartera para acciones de retencion.",
        "Valor: usar datos de contrato, servicios y facturacion para llegar antes.",
        "Trade-off: priorizamos recall, pero monitoreamos precision para no sobreactuar."
    ], 0.7, 1.45, 5.2, 2.25, 14)
    card(s, "SALIDA", "score", "probabilidad de churn", 6.35, 1.45, 2.6, 0.95, BLUE)
    card(s, "USO", "ranking", "agenda de retencion", 6.35, 2.65, 2.6, 0.95, GREEN)
    card(s, "RIESGO", "falsos +", "costo comercial", 6.35, 3.85, 2.6, 0.95, RED)

    # 3 dataset
    s = new("2 · Dataset", "7.043 clientes, 21 variables y una clase positiva minoritaria", 3)
    card(s, "TARGET", "Churn", "Yes / No", 0.65, 1.35, 1.8, 1.0, ORANGE)
    card(s, "POSITIVOS", pct(df["churn_num"].mean()), "churn real", 2.75, 1.35, 1.8, 1.0, RED)
    card(s, "LIMPIEZA", "11", "TotalCharges vacio, tenure=0", 4.85, 1.35, 2.1, 1.0, TEAL)
    card(s, "FEATURES", "originales + FE", "sin target encoding", 7.25, 1.35, 2.0, 1.0, BLUE)
    bullets(s, [
        "Variables demograficas: genero, senior, pareja, dependientes.",
        "Servicios: internet, seguridad, soporte, backup, streaming.",
        "Contrato/facturacion: tenure, tipo de contrato, metodo de pago, cargos.",
        "Limitacion estructural: no hay fecha de corte ni historial temporal."
    ], 0.8, 2.95, 8.5, 1.45, 12.5)

    # 4 EDA
    s = new("3 · EDA", "La tasa de churn cambia mucho segun contrato, servicio y soporte", 4)
    image(s, ASSETS / "eda_churn_rate_variables.png", 0.55, 1.25, 8.95)

    # 5 EDA 2
    s = new("3 · EDA", "El patron principal: clientes nuevos y poco anclados", 5)
    image(s, ASSETS / "eda_target_tenure.png", 0.6, 1.3, 5.0)
    bullets(s, [
        "Contrato mensual y baja antiguedad concentran mayor riesgo.",
        "Fibra optica aparece como señal predictiva relevante.",
        "Servicios de soporte/seguridad reducen riesgo observado.",
        "Estas relaciones son asociaciones predictivas, no causalidad."
    ], 6.05, 1.55, 3.25, 2.5, 12.5)

    # 6 evaluacion
    s = new("4 · Evaluacion", "PR-AUC para elegir modelo; F2 para elegir umbral", 6)
    card(s, "TRAIN", "5.634", "80%, estratificado", 0.65, 1.35, 1.8, 1.0, BLUE)
    card(s, "TEST", "1.409", "20%, cerrado", 2.75, 1.35, 1.8, 1.0, TEAL)
    card(s, "CV", "3 folds", "solo train", 4.85, 1.35, 1.8, 1.0, GREEN)
    card(s, "NO SKILL", num(df["churn_num"].mean()), "PR-AUC base", 6.95, 1.35, 1.8, 1.0, RED)
    bullets(s, [
        "PR-AUC mide la calidad del ranking en una clase positiva minoritaria.",
        "F2 prioriza recall para el punto operativo, sin ignorar precision.",
        "El test no se usa para tuning, seleccion de modelo ni seleccion de umbral.",
        "Pipeline: imputacion, escalado y one-hot se ajustan dentro de cada fold."
    ], 0.85, 3.0, 8.4, 1.4, 12.5)

    # 7 baseline
    s = new("5 · Baseline", "La regresion logistica deja una vara alta e interpretable", 7)
    cvi = cv.set_index("model")
    card(s, "DUMMY PR-AUC", num(cvi.loc["dummy_most_frequent", "pr_auc_mean"]), "sin informacion", 0.8, 1.45, 2.0, 1.0, RED)
    card(s, "LOGISTIC PR-AUC", num(cvi.loc["logistic_regression", "pr_auc_mean"]), "baseline fuerte", 3.1, 1.45, 2.0, 1.0, BLUE)
    card(s, "LOGISTIC F2", num(cvi.loc["logistic_regression", "f2_mean"]), "recall alto", 5.4, 1.45, 2.0, 1.0, GREEN)
    bullets(s, [
        "La logistica queda cerca de CatBoost: el problema tiene señal lineal fuerte.",
        "Se usa como baseline interpretable para exigir que modelos complejos justifiquen su mejora.",
        "El dummy confirma que accuracy no alcanza: predecir siempre No Churn no sirve para retener."
    ], 1.05, 3.05, 7.9, 1.25, 13)

    # 8 seleccion
    s = new("6 · Seleccion de modelos", "10 pruebas comparadas por PR-AUC", 8)
    image(s, ASSETS / "modelos_10_prauc.png", 0.55, 1.22, 6.45)
    top = cv.sort_values("pr_auc_mean", ascending=False).head(3)
    y = 1.45
    for _, r in top.iterrows():
        card(s, r["model"], num(r["pr_auc_mean"]), f"F2 {num(r['f2_mean'])} · Precision {num(r['precision_mean'])}", 7.25, y, 2.25, 0.78, ORANGE if r["model"] == metrics["selected_model"] else BLUE)
        y += 0.95
    textbox(s, "CatBoost gana por poco. Logistic y Gradient Boosting quedan cerca, por eso se reporta el ranking completo y no solo el ganador.", 7.25, 4.3, 2.25, 0.7, 9.8, MUTED)

    # 9 modelo final
    s = new("7 · Modelo final", "CatBoost: mejor PR-AUC y umbral F2 = 0.31", 9)
    card(s, "PR-AUC", num(metrics["test_pr_auc"]), "test", 0.6, 1.35, 1.55, 0.95, ORANGE)
    card(s, "ROC-AUC", num(metrics["test_roc_auc"]), "test", 2.35, 1.35, 1.55, 0.95, BLUE)
    card(s, "Recall", pct(metrics["test_recall"]), "detecta churn", 4.1, 1.35, 1.55, 0.95, GREEN)
    card(s, "Precision", pct(metrics["test_precision"]), "contactos utiles", 5.85, 1.35, 1.55, 0.95, RED)
    card(s, "F2", num(metrics["test_f2"]), "umbral", 7.6, 1.35, 1.55, 0.95, TEAL)
    image(s, ASSETS / "threshold_precision_recall.png", 0.75, 2.75, 5.2)
    image(s, ASSETS / "capture_curve.png", 5.95, 2.75, 3.45)

    # 10 interpretacion
    s = new("7 · Modelo final", "Interpretacion: contrato, tenure y señales de servicio", 10)
    image(s, ASSETS / "importance_permutation.png", 0.55, 1.25, 4.65)
    image(s, ASSETS / "importance_shap.png", 5.05, 1.25, 4.55)
    textbox(s, "Permutation importance esta alineada a F2; SHAP resume impacto global. No es causalidad: es señal predictiva para priorizar.", 0.85, 4.85, 8.2, 0.35, 10.5, MUTED, align=PP_ALIGN.CENTER)

    # 11 limitaciones
    s = new("8 · Limitaciones / mejoras", "La mayor mejora no es otro algoritmo: son datos y decision", 11)
    bullets(s, [
        "No hay corte temporal real: validar con meses futuros seria la mejora metodologica mas importante.",
        "Precision moderada: el umbral debe ajustarse con costos reales y capacidad comercial.",
        "Variables redundantes: algunas originales y derivadas comparten informacion.",
        "Agregar reclamos, uso real, tickets, promociones, cambios de plan y satisfaccion.",
        "Uplift modeling: contactar clientes persuadibles, no solo clientes con alto riesgo."
    ], 0.8, 1.45, 8.3, 2.5, 13)
    card(s, "LEAKAGE", "sin señal grave", "ID fuera, target fuera, pipeline por fold", 0.9, 4.25, 2.5, 0.85, GREEN)
    card(s, "OVERFIT", "controlado", "gap F2 bajo/moderado", 3.75, 4.25, 2.5, 0.85, BLUE)
    card(s, "PROXIMO PASO", "costos", "umbral por valor esperado", 6.6, 4.25, 2.5, 0.85, ORANGE)

    # 12 conclusiones
    s = new("9 · Conclusiones", "Un ranking accionable para retencion", 12)
    card(s, "METODO", "PR-AUC + F2", "ranking y decision", 0.75, 1.35, 2.35, 1.05, BLUE)
    card(s, "MODELO", "CatBoost", "mejor CV", 3.4, 1.35, 2.35, 1.05, ORANGE)
    card(s, "NEGOCIO", "priorizar", "agenda de retencion", 6.05, 1.35, 2.35, 1.05, GREEN)
    bullets(s, [
        "El churn se concentra en perfiles con contrato mensual, baja antiguedad y menor anclaje de servicios.",
        "CatBoost fue el mejor modelo por PR-AUC dentro de 10 pruebas comparables.",
        "El umbral actual prioriza recall; debe calibrarse con costos y capacidad real.",
        "El modelo es una herramienta de priorizacion, no una explicacion causal del abandono."
    ], 0.95, 3.0, 8.1, 1.45, 13.5)

    prs.save(PPTX)


def main():
    df, metrics, cv, thresholds, preds, perm, shap, checks = load()
    make_assets(df, metrics, cv, thresholds, preds, perm, shap)
    build_deck(df, metrics, cv, preds, checks)
    print(PPTX)


if __name__ == "__main__":
    main()
